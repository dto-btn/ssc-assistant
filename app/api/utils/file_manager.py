from pypdf import PdfReader
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
from io import BytesIO
import csv
from typing import Iterable

class FileManager:
    def __init__(self, data: bytes, filetype: str):
        self.data = data
        self.filetype = filetype
        self.content = None

    def extract_text(self) -> str:
        # Normalize filetype for easier matching
        ft = (self.filetype or '').lower()

        # PDF
        if ft == 'application/pdf':
            return self._extract_pdf(self.data)

        # Word (docx) and common legacy/macro-enabled MIME for .doc (best-effort)
        if ft in {
            'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
            'application/msword',  # legacy .doc (will likely fail with python-docx)
            'application/vnd.ms-word.document.macroenabled.12',
        }:
            try:
                return self._extract_docx(self.data)
            except Exception:
                return "[Unsupported Word format. Please upload .docx]"

        # Excel (xlsx) and common legacy/macro-enabled MIME
        if ft in {
            'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
            'application/vnd.ms-excel',  # legacy .xls (openpyxl won't read .xls)
            'application/vnd.ms-excel.sheet.macroenabled.12',
        }:
            try:
                return self._extract_xlsx(self.data)
            except Exception:
                return "[Unsupported Excel format. Please upload .xlsx]"

        # CSV / TSV
        if ft in {'text/csv', 'application/csv', 'text/tab-separated-values'}:
            return self._extract_csv(self.data)

        # PowerPoint (pptx)
        if ft in {
            'application/vnd.openxmlformats-officedocument.presentationml.presentation',
            'application/vnd.ms-powerpoint',  # legacy .ppt
            'application/vnd.ms-powerpoint.presentation.macroenabled.12',
        }:
            try:
                return self._extract_pptx(self.data)
            except Exception:
                return "[Unsupported PowerPoint format. Please upload .pptx]"

        # Plain text
        if ft == 'text/plain' or ft.startswith('text/'):
            return self._extract_txt(self.data)

        # Fallback
        return "[Unsupported file type]"

    def _extract_pdf(self, data: bytes) -> str:
        try:
            reader = PdfReader(BytesIO(data))
            text = "\n".join(page.extract_text() or '' for page in reader.pages)
            return text
        except Exception:
            return "[Failed to read PDF]"

    def _extract_docx(self, data: bytes) -> str:
        doc = Document(BytesIO(data))
        text = "\n".join([p.text for p in doc.paragraphs])
        return text

    def _extract_xlsx(self, data: bytes) -> str:
        wb = load_workbook(BytesIO(data), read_only=True)
        text = []
        for ws in wb.worksheets:
            for row in ws.iter_rows(values_only=True):
                text.append("\t".join([str(cell) if cell is not None else '' for cell in row]))
        return "\n".join(text)

    def _extract_pptx(self, data: bytes) -> str:
        prs = Presentation(BytesIO(data))
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)  # type: ignore
        return "\n".join(text)

    def _extract_txt(self, data: bytes) -> str:
        # Try a few common encodings, then replace errors to avoid crashes
        for enc in ('utf-8', 'utf-8-sig', 'utf-16', 'utf-16le', 'utf-16be', 'latin-1'):
            try:
                return data.decode(enc)
            except UnicodeDecodeError:
                continue
        return data.decode('utf-8', errors='replace')

    def _extract_csv(self, data: bytes) -> str:
        # Decode CSV robustly and join cells with tabs
        decoded = None
        for enc in ('utf-8', 'utf-8-sig', 'latin-1'):
            try:
                decoded = data.decode(enc)
                break
            except UnicodeDecodeError:
                continue
        if decoded is None:
            decoded = data.decode('utf-8', errors='replace')

        reader = csv.reader(decoded.splitlines())
        lines: Iterable[str] = ("\t".join(row) for row in reader)
        return "\n".join(lines)

